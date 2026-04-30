import io
import pytest
from ext.services.extractor import extract_text, UnsupportedMimeType


def test_txt_passthrough():
    assert extract_text(b"hello", "text/plain", "a.txt") == "hello"


def test_markdown_passthrough():
    assert extract_text(b"# Title\n\npara", "text/markdown", "a.md") == "# Title\n\npara"


def test_unsupported_mime_raises():
    with pytest.raises(UnsupportedMimeType):
        extract_text(b"...", "application/vnd.ms-excel", "a.xls")


def test_extractor_dispatches_on_extension_if_mime_missing():
    assert extract_text(b"hi", "application/octet-stream", "note.txt") == "hi"


def test_pdf_extraction_tiny():
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(width=612, height=792)
    buf = io.BytesIO()
    w.write(buf)
    assert extract_text(buf.getvalue(), "application/pdf", "a.pdf") == ""


# ---------------------------------------------------------------------------
# PPTX support (Task 2 — slide-level extraction with python-pptx).
# ---------------------------------------------------------------------------

def _make_pptx() -> bytes:
    """Build a tiny .pptx in memory with title + body + table + notes."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1: title + body bullets + speaker notes.
    s1 = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    s1.shapes.title.text = "Quarterly Review"
    s1.placeholders[1].text = "Revenue up 20% YoY\nCustomer count: 500"
    s1.notes_slide.notes_text_frame.text = "Highlight churn risk"

    # Slide 2: title + table.
    s2 = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    s2.shapes.title.text = "Pipeline"
    table = s2.shapes.add_table(
        rows=2, cols=2,
        left=Inches(1), top=Inches(2),
        width=Inches(4), height=Inches(1.5),
    ).table
    table.cell(0, 0).text = "Stage"
    table.cell(0, 1).text = "Count"
    table.cell(1, 0).text = "Qualified"
    table.cell(1, 1).text = "42"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_extraction_text_roundtrip():
    """Flat extractor surfaces title, body, table cells, and speaker notes."""
    data = _make_pptx()
    text = extract_text(
        data,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "deck.pptx",
    )
    for needle in [
        "Quarterly Review",
        "Revenue up 20% YoY",
        "Customer count: 500",
        "Highlight churn risk",
        "Pipeline",
        "Stage",
        "Count",
        "Qualified",
        "42",
    ]:
        assert needle in text, f"missing {needle!r} in extracted text"


def test_pptx_extraction_resolves_by_extension():
    """When MIME is unknown, .pptx extension drives the resolver."""
    data = _make_pptx()
    text = extract_text(data, "application/octet-stream", "deck.pptx")
    assert "Quarterly Review" in text


def test_pptx_blocks_emit_per_slide_with_page_numbers():
    """Structural extractor returns one ExtractedBlock per slide with page set."""
    pytest.importorskip("pptx")
    from ext.services.extractor import extract

    data = _make_pptx()
    blocks = extract(
        data,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "deck.pptx",
    )
    assert len(blocks) == 2, f"expected 2 slides, got {len(blocks)}"
    assert blocks[0].page == 1
    assert blocks[1].page == 2
    # Title placeholder lands in heading_path on slide 1.
    assert blocks[0].heading_path == ["Quarterly Review"]
    # Speaker notes flow into the body text.
    assert "Highlight churn risk" in blocks[0].text
    # Table cells flow into slide 2's body.
    assert "Stage" in blocks[1].text
    assert "Qualified" in blocks[1].text
    assert "42" in blocks[1].text
