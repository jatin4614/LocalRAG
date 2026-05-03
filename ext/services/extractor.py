"""Text extraction from uploaded documents.

Two entry points:

* :func:`extract` returns ``list[ExtractedBlock]`` with structural metadata
  (page, heading_path, sheet). Preferred for all new callers — downstream
  ingest uses the metadata to enrich Qdrant payloads so retrieval can
  surface "see page 7" / "under heading Foo > Bar" hints.

* :func:`extract_flat` (and the legacy :func:`extract_text` alias) return a
  single flat ``str``, byte-identical to the pre-P0.4 implementation. Old
  callers that only want text are unaffected.

The old per-format ``_extract_*`` helpers are preserved verbatim and power
``extract_flat`` so that legacy string output does not regress. The new
structural extractors live in ``_blocks_*`` helpers.
"""
from __future__ import annotations

import io
import logging
import os
import re
from dataclasses import dataclass, field
from typing import Callable, Optional


log = logging.getLogger(__name__)


class UnsupportedMimeType(RuntimeError):
    pass


# ---------------------------------------------------------------------------
# Structural block
# ---------------------------------------------------------------------------

@dataclass
class ExtractedBlock:
    """A chunk of extracted text together with its structural origin.

    Fields that a given format can't supply stay at their default
    (``None`` / ``[]``). Consumers should treat missing metadata as "unknown"
    rather than inventing values.

    ``kind`` is consumed by ``ext.services.ingest._coalesce_small_blocks`` to
    decide which blocks may merge. Defaults to ``"prose"`` (the common case).
    Extractors that emit non-prose structural blocks (e.g., docx tables → TSV
    rows; xlsx sheets → TSV rows) stamp ``kind="table"`` so the coalescer
    leaves them atomic. Adding a new value here is non-breaking — the
    coalescer treats anything other than ``"prose"`` as atomic.
    """

    text: str
    page: Optional[int] = None  # 1-based PDF page number; None for non-paginated types.
    heading_path: list[str] = field(default_factory=list)  # DOCX/MD heading stack, root -> leaf.
    sheet: Optional[str] = None  # XLSX worksheet name; None otherwise.
    kind: str = "prose"  # "prose" | "table" | future: "code" | "image_caption"


# ---------------------------------------------------------------------------
# Legacy flat-string extractors — kept byte-identical to pre-P0.4.
# ---------------------------------------------------------------------------

def _extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []

    def _iter_block_items(parent):
        # Walk body elements in document order so prose and tables interleave correctly.
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        body = parent.element.body if hasattr(parent, "element") else parent._element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, parent)
            elif child.tag == qn("w:tbl"):
                yield Table(child, parent)

    for block in _iter_block_items(doc):
        if block.__class__.__name__ == "Paragraph":
            txt = block.text.strip()
            if txt:
                parts.append(txt)
        else:  # Table
            for row in block.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append("\t".join(cells))

    # Headers and footers (policy docs often put key info here).
    for section in doc.sections:
        for container in (section.header, section.footer):
            for para in container.paragraphs:
                txt = para.text.strip()
                if txt:
                    parts.append(txt)

    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_csv(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_pptx(data: bytes) -> str:
    """Flat-string extractor for .pptx — slide-by-slide text concatenation.

    Per slide: title placeholder + body text + table cells + speaker notes.
    Slides joined with blank-line breaks so downstream tokenizers see slide
    boundaries. Speaker notes carry the analytical content the slide image
    only hints at, so they're included verbatim.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        slide_parts.append("\t".join(cells))
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                slide_parts.append(notes)
        if slide_parts:
            parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


EXTRACTORS: dict[str, Callable[[bytes], str]] = {
    "text/plain": _extract_txt,
    "text/markdown": _extract_txt,
    "text/csv": _extract_csv,
    "application/pdf": _extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
}


_EXT_FALLBACK = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".csv": "text/csv",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _resolve_mime(mime_type: str, filename: str) -> Optional[str]:
    """Return the registered mime key for ``mime_type`` or filename fallback,
    or None when nothing matches. Raises on legacy .doc early.
    """
    if filename.lower().endswith(".doc") and not filename.lower().endswith(".docx"):
        raise UnsupportedMimeType(
            f"legacy .doc format not supported (convert to .docx first): {filename}"
        )
    if mime_type in EXTRACTORS:
        return mime_type
    for ext, m in _EXT_FALLBACK.items():
        if filename.lower().endswith(ext):
            return m
    return None


# ---------------------------------------------------------------------------
# Structural block extractors — new in P0.4.
# ---------------------------------------------------------------------------

def _blocks_txt(data: bytes) -> list[ExtractedBlock]:
    # Plain text has no structure; emit a single block. Empty input → no blocks.
    text = _extract_txt(data)
    if not text:
        return []
    return [ExtractedBlock(text=text)]


def _blocks_csv(data: bytes) -> list[ExtractedBlock]:
    text = _extract_csv(data)
    if not text:
        return []
    return [ExtractedBlock(text=text)]


def _blocks_pdf(data: bytes) -> list[ExtractedBlock]:
    """Extract text per page. Prefer pymupdf when available — it deduplicates
    overlay/watermark text streams that pypdf returns verbatim, which on
    stamped/scanned policy PDFs (e.g. ARMY CYBER SECURITY POLICY 2023) caused
    the same page text to appear twice in every chunk (~59% chunks duplicated,
    measured 2026-04-28). pypdf path retained as the import-failure fallback.
    """
    try:
        import pymupdf  # type: ignore
        doc = pymupdf.open(stream=data, filetype="pdf")
        try:
            blocks: list[ExtractedBlock] = []
            for i in range(doc.page_count):
                page = doc.load_page(i)
                t = page.get_text() or ""
                t = _dedup_overlay_text(t)
                if not t.strip():
                    continue
                blocks.append(ExtractedBlock(text=t, page=i + 1))
            return blocks
        finally:
            doc.close()
    except ImportError:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        blocks = []
        for i, page in enumerate(reader.pages, start=1):
            t = page.extract_text() or ""
            t = _dedup_overlay_text(t)
            if not t.strip():
                continue
            blocks.append(ExtractedBlock(text=t, page=i))
        return blocks


_HEADING_RE = re.compile(r"^Heading\s+([1-9])$", re.IGNORECASE)


def _dedup_overlay_text(text: str) -> str:
    """Strip duplicate-content artifacts that some stamped/watermarked PDFs
    embed per page. Observed pattern: ``A + W + A`` where ``A`` is the real
    page text and ``W`` is a one-line stamp (e.g. ``IP / TIMESTAMP``); both
    pypdf and pymupdf faithfully extract the doubled content.

    Algorithm: take the first 80 stripped chars as a fingerprint; search for
    them past the first third of ``text``; if char-match between the leading
    half and the candidate duplicate exceeds 85%, truncate at the duplicate
    boundary, returning ``A + W`` (the stamp is harmless context). Returns
    input unchanged when no duplication found — conservative; won't fire on
    pages with merely overlapping vocabulary or short repeated phrases.

    Set ``RAG_PDF_DEDUP=0`` to disable. Default on. (Plan B Phase 6 followup
    2026-04-28 — kb_3 ARMY CYBER SECURITY POLICY 2023.pdf showed 58 % chunks
    duplicated; this helper drops that to ~0 %.)
    """
    import os
    if os.environ.get("RAG_PDF_DEDUP", "1") != "1":
        return text
    if not text or len(text) < 200:
        return text
    stripped = text.lstrip()
    if len(stripped) < 200:
        return text
    head_len = min(80, len(stripped) // 3)
    if head_len < 30:
        return text
    head = stripped[:head_len]
    second = text.find(head, len(text) // 3)
    if second < 0:
        return text
    a, b = text[:second], text[second:]
    # Skip if the fingerprint reappears past its own leading occurrence in b —
    # that's a 3+x repeat (legitimate content like a TOC or chorus), not the
    # binary A+W+A artifact we're trying to fix.
    if head in b[len(head):]:
        return text
    common = min(len(a), len(b))
    if common < 30:
        return text
    matches = sum(1 for x, y in zip(a[:common], b[:common]) if x == y)
    if matches / common > 0.85:
        return text[:second].rstrip()
    return text


def _blocks_docx(data: bytes) -> list[ExtractedBlock]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    heading_stack: list[str] = []
    blocks: list[ExtractedBlock] = []

    # 1) Paragraphs: track heading levels via style.name == "Heading N".
    for para in doc.paragraphs:
        txt = para.text
        if not txt.strip():
            continue
        style_name = getattr(getattr(para, "style", None), "name", None) or ""
        m = _HEADING_RE.match(style_name.strip())
        if m:
            level = int(m.group(1))
            # Pop deeper levels, then set this level.
            del heading_stack[level - 1 :]
            heading_stack.append(txt.strip())
            # Headings themselves are not emitted as prose blocks — they only
            # update the stack. (Keeps parity with legacy docx behaviour that
            # included them as regular paragraphs; flat extractor still does.)
            continue
        blocks.append(
            ExtractedBlock(text=txt, heading_path=list(heading_stack))
        )

    # 2) Tables: TSV flatten, one block per table, carrying the current
    # heading_path at the point the table appears. python-docx does not give
    # a single ordered iterator over paragraphs-and-tables, so tables are
    # emitted after all paragraphs. That differs from the legacy flat output
    # (which omits tables entirely — only paragraphs were joined); kept here
    # as additive structural info.
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            blocks.append(
                ExtractedBlock(
                    text="\n".join(rows),
                    heading_path=list(heading_stack),
                    kind="table",
                )
            )

    return blocks


def _blocks_xlsx(data: bytes) -> list[ExtractedBlock]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[ExtractedBlock] = []
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if not rows:
            # Skip empty sheets entirely; keeps the output aligned with
            # legacy flat behaviour which produced no bytes for such sheets.
            continue
        blocks.append(ExtractedBlock(
            text="\n".join(rows), sheet=ws.title, kind="table",
        ))
    return blocks


def _blocks_pptx(data: bytes) -> list[ExtractedBlock]:
    """Structural extractor for .pptx — one block per slide.

    Per slide: title placeholder (when present) becomes the heading_path;
    body text + table rows + speaker notes accumulate into the slide's
    text. The slide number lives in ``page`` (analogous to PDF page
    numbers) so retrieval can render "see slide N" hints. Empty slides
    are skipped.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    blocks: list[ExtractedBlock] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        title: Optional[str] = None
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_txt: list[str] = []
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        shape_txt.append(txt)
                if not shape_txt:
                    continue
                # Title placeholder (idx == 0) becomes the slide heading.
                if title is None and getattr(shape, "is_placeholder", False):
                    try:
                        ph_idx = shape.placeholder_format.idx
                    except Exception:
                        ph_idx = None
                    if ph_idx == 0:
                        title = "\n".join(shape_txt).strip()
                        continue
                body_parts.append("\n".join(shape_txt))
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        body_parts.append("\t".join(cells))
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                body_parts.append(notes)
        body = "\n\n".join(p for p in body_parts if p).strip()
        if not body and not title:
            continue
        text = body if body else (title or "")
        heading_path = [title] if title else []
        blocks.append(
            ExtractedBlock(text=text, page=slide_idx, heading_path=heading_path)
        )
    return blocks


_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
_MD_PARA_SPLIT = re.compile(r"\n\s*\n")


def _blocks_md(data: bytes) -> list[ExtractedBlock]:
    # Markdown: track heading stack, split body into paragraphs per current
    # heading_path. Setext-style headings (``Title\n====``) are NOT handled —
    # keep the parser simple; ATX-style covers the common case.
    text = _extract_txt(data)
    if not text:
        return []
    heading_stack: list[str] = []
    blocks: list[ExtractedBlock] = []
    # Work line-by-line to keep heading tracking precise, then re-group
    # consecutive non-heading lines into paragraphs.
    buf: list[str] = []

    def flush() -> None:
        if not buf:
            return
        joined = "\n".join(buf).strip()
        buf.clear()
        if not joined:
            return
        # Split on blank-line boundaries within the buffer — preserves
        # paragraph structure under the current heading.
        for para in _MD_PARA_SPLIT.split(joined):
            p = para.strip()
            if p:
                blocks.append(
                    ExtractedBlock(text=p, heading_path=list(heading_stack))
                )

    for line in text.splitlines():
        m = _MD_HEADING_RE.match(line)
        if m:
            flush()
            level = len(m.group(1))
            title = m.group(2).strip()
            del heading_stack[level - 1 :]
            heading_stack.append(title)
        else:
            buf.append(line)
    flush()
    return blocks


# ---------------------------------------------------------------------------
# Optional layout-aware extraction via the ``unstructured`` library.
# ---------------------------------------------------------------------------

def _blocks_unstructured(
    data: bytes, mime: str, filename: str,
) -> list[ExtractedBlock]:
    """Layout-aware element extraction via ``unstructured.partition.auto``.

    Opt-in via ``RAG_UNSTRUCTURED_LAYOUT=1``. When enabled, replaces the
    legacy per-page (PDF) or per-slide (PPTX) extractors with a typed-
    element parse — Title / Header / Text / NarrativeText / Table /
    ListItem — that preserves reading order across multi-column layouts
    and surfaces table structure as element-native metadata.

    ``RAG_UNSTRUCTURED_STRATEGY`` selects the parser strategy:
      * ``fast`` (default) — heuristic extraction, no extras required.
      * ``hi_res`` — Detectron2 + LayoutParser layout reconstruction;
        operator must install the ``unstructured[local-inference]`` extras
        in the open-webui image (celery worker is intentionally lean and
        does not carry the heavyweight layout deps).
      * ``ocr_only`` — Tesseract-based OCR for scanned docs.

    Mapping to ExtractedBlock:
      * ``Title`` / ``Header`` → push onto a single-level heading stack
        (we don't try to infer depth — unstructured doesn't expose it
        reliably across formats).
      * Everything else → emit as a block with the current heading_path
        and the element's ``page_number`` (when present).

    Raises ``ImportError`` if the library is missing, so callers can
    fall back to the legacy extractor.
    """
    from unstructured.partition.auto import partition

    strategy = os.environ.get("RAG_UNSTRUCTURED_STRATEGY", "fast")
    elements = partition(
        file=io.BytesIO(data),
        content_type=mime,
        metadata_filename=filename,
        strategy=strategy,
    )

    blocks: list[ExtractedBlock] = []
    heading_stack: list[str] = []
    for el in elements:
        text = (getattr(el, "text", None) or "").strip()
        if not text:
            continue
        meta = getattr(el, "metadata", None)
        page = getattr(meta, "page_number", None) if meta else None
        cat = (
            getattr(el, "category", None)
            or el.__class__.__name__
        )
        if cat in ("Title", "Header"):
            heading_stack = [text]
            continue
        blocks.append(ExtractedBlock(
            text=text,
            page=page,
            heading_path=list(heading_stack),
        ))
    return blocks


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract(data: bytes, mime_type: str, filename: str) -> list[ExtractedBlock]:
    """Return structural blocks for ``data``. Dispatches on ``mime_type`` with
    filename-extension fallback. Raises :class:`UnsupportedMimeType` for
    unknown types (and legacy .doc binaries).

    Opt-in: when ``RAG_UNSTRUCTURED_LAYOUT=1`` and ``resolved`` is PDF or
    PPTX, prefer ``unstructured.partition.auto`` for layout-aware
    extraction. Falls back to the legacy per-format extractor on
    ImportError (library not installed) or any partition failure.
    """
    resolved = _resolve_mime(mime_type, filename)
    if resolved is None:
        raise UnsupportedMimeType(f"{mime_type} (filename={filename})")
    if (
        os.environ.get("RAG_UNSTRUCTURED_LAYOUT", "0") == "1"
        and resolved in (
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    ):
        try:
            return _blocks_unstructured(data, resolved, filename)
        except ImportError:
            pass
        except Exception as e:  # noqa: BLE001 — fail-open to legacy
            log.warning(
                "unstructured extraction failed for %s: %s; falling back",
                filename, e,
            )
    if resolved == "text/markdown":
        return _blocks_md(data)
    if resolved == "text/plain":
        return _blocks_txt(data)
    if resolved == "text/csv":
        return _blocks_csv(data)
    if resolved == "application/pdf":
        return _blocks_pdf(data)
    if resolved == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _blocks_docx(data)
    if resolved == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return _blocks_xlsx(data)
    if resolved == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _blocks_pptx(data)
    # _resolve_mime only returns keys present in EXTRACTORS — unreachable.
    raise UnsupportedMimeType(f"{mime_type} (filename={filename})")


def extract_flat(data: bytes, mime_type: str, filename: str) -> str:
    """Back-compat flat-string extractor. Byte-identical to the pre-P0.4
    implementation for all supported formats."""
    resolved = _resolve_mime(mime_type, filename)
    if resolved is None:
        raise UnsupportedMimeType(f"{mime_type} (filename={filename})")
    fn = EXTRACTORS[resolved]
    return fn(data)


# Legacy alias — existing imports (``from ext.services.extractor import
# extract_text``) keep working.
extract_text = extract_flat
